# coding:utf-8
import json
from pptx import Presentation
from pptx.dml.color import RGBColor
from collections import OrderedDict
from pptx.util import Pt #Inches

class Slide():
  def __init__(self,prs,slide_source):
    self.__shapes = prs.slides.add_slide(prs.slide_layouts[slide_source.get('layout_number','1')]).shapes
    #placeholders[0] is title, placeholders[1] is maintext
    self.__make_title(slide_source)

    #ページ内に入れたい文章一つにつき1回
    if 'contexts' in slide_source:
      for paragraph_number in slide_source.get('contexts','null'):
        self.__addpara(
          self.__shapes.placeholders[1].text_frame.add_paragraph(),
          slide_source['contexts'][str(paragraph_number)]
        )
  #新しくパラグラフを追加(文章1つ分)
  def __addpara(self,this_paragraph,para_source):
    this_paragraph.text = para_source.get('text','無')
    this_paragraph.level = para_source.get('level',0)
    if 'size' in para_source:
      this_paragraph.font.size = Pt(para_source['size'])
    if 'color' in para_source:
      this_paragraph.font.color.rgb = RGBColor.from_string(para_source['color'])


  def __make_title(self,slide_source):
    self.__addpara(
      self.__shapes.placeholders[0].text_frame.add_paragraph(),
      slide_source.get('title','')
    )
