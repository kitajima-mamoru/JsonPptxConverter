# coding:utf-8
import json
from pptx import Presentation
from pptx.dml.color import RGBColor
from collections import OrderedDict
from pptx.util import Pt #Inches

class getdata():
  def __init__(self,params):
    #jsonをロード
    self.master_sourse = json.load(open(params['sourse'],encoding = "utf-8"))
    #新規pptxファイルの名前
    self.presentation_name = self.master_sourse['0']['presentation_name']
    #レイアウトをコピーし別名で保存　将来的には設定用ディレクトリを作りそちらに格納する
    Presentation(params['template']).save(self.presentation_name)
    self.prs = Presentation(self.presentation_name)
    self.slide_sum = params['slide_sum']
  #新しくpptxを作成
  def main(self):
    #各ページを追加、保存していく
    for slide_number in range(1,self.slide_sum+1):
      self.__addslide(self.master_sourse[str(slide_number)])

  #新しくスライド追加(スライド1つ分)
  def __addslide(self,slide_sourse):
    #slidemasterからlayoutを指定しつつslideを追加
    shapes = self.prs.slides.add_slide(self.prs.slide_layouts[slide_sourse['layout_number']]).shapes
    #placeholders[0] is title, placeholders[1] is maintext
    self.__addpara(
        shapes.placeholders[0].text_frame.add_paragraph(),
        slide_sourse.get('title','')
    )
    #ページ内に入れたい文章一つにつき1回
    if 'contexts' in slide_sourse:
      for paragraph_number in slide_sourse.get('contexts','null'):
        self.__addpara(
            shapes.placeholders[1].text_frame.add_paragraph(),
            slide_sourse['contexts'][str(paragraph_number)]
        )
  #新しくパラグラフを追加(文章1つ分)
  def __addpara(self,this_paragraph,para_sourse):
    this_paragraph.text = para_sourse.get('text','無')
    this_paragraph.level = para_sourse.get('level',0)
    if 'size' in para_sourse:
      this_paragraph.font.size = Pt(para_sourse['size'])
    if 'color' in para_sourse:
      this_paragraph.font.color.rgb = RGBColor.from_string(para_sourse['color'])

  def save(self):
    self.prs.save(self.presentation_name)
